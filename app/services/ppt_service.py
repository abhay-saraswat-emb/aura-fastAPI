import os
from openai import OpenAI
import anthropic
import google.generativeai as genai
from dotenv import load_dotenv
from typing import Union, AsyncGenerator, Any
from google import genai
from google.genai import types
from typing import List, Dict, Optional
import httpx
from io import BytesIO
from pptx import Presentation

load_dotenv()

openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
claude_client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
genai_client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))


async def ppt_claude_stream_response(
    prompt: str, 
    system_prompt: str, 
    max_tokens: int, 
    temperature: float,
    conversation_history: List[Dict[str, str]] = [],
    ppt_url: str = None
):
    """
    Stream a response from Claude model for PPT content.
    
    Args:
        prompt: The user's input prompt
        system_prompt: Instructions for the AI
        max_tokens: Maximum response length
        temperature: Randomness of response (0-1)
        conversation_history: history of conversation
        ppt_url: URL of PowerPoint file
        
    Yields:
        Chunks of the generated response
    """
    try:
        # Convert conversation history to Claude format
        messages = [{"role": msg["role"], "content": msg["content"]} for msg in conversation_history]
        user_message = []

        # Extract and format text from the PPT file
        if ppt_url:
            try:
                response = httpx.get(ppt_url)
                response.raise_for_status()

                prs = Presentation(BytesIO(response.content))
                ppt_text = []

                for i, slide in enumerate(prs.slides):
                    slide_text = [f"--- Slide {i + 1} ---"]
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text.strip())
                    ppt_text.append("\n".join(slide_text))

                full_ppt_text = "\n\n".join(ppt_text)

                # Add extracted PPT content
                user_message.append({
                    "type": "text",
                    "text": f"Here is the content of the PowerPoint presentation:\n{full_ppt_text}"
                })

            except Exception as e:
                yield f"data: Error extracting PPT: {str(e)}\n\n"
                return

        # Add the user prompt
        user_message.append({
            "type": "text",
            "text": prompt
        })

        messages.append({"role": "user", "content": user_message})

        # Stream response from Claude
        with claude_client.messages.stream(
            model="claude-3-7-sonnet-20250219",
            system=system_prompt,
            messages=messages,
            max_tokens=max_tokens,
            temperature=temperature
        ) as stream:
            for text in stream.text_stream:
                yield text

    except Exception as e:
        yield f"data: Error: {str(e)}\n\n"

async def ppt_genai_stream_response(
    prompt: str, 
    system_prompt: str, 
    max_tokens: int, 
    temperature: float,
    conversation_history: List[Dict[str, str]] = [],
    ppt_url: str = None
):
    try:
        # Format conversation history for Gemini
        history_content = [
            types.Content(
                role="user" if msg["role"] == "user" else "model",
                parts=[types.Part(text=msg["content"])]
            ) for msg in conversation_history
        ]

        # Extract PPTX content
        if ppt_url:
            try:
                response = httpx.get(ppt_url)
                response.raise_for_status()

                prs = Presentation(BytesIO(response.content))
                ppt_text = []

                for i, slide in enumerate(prs.slides):
                    slide_text = [f"--- Slide {i + 1} ---"]
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text.strip())
                    ppt_text.append("\n".join(slide_text))

                full_ppt_text = "\n\n".join(ppt_text)

                # Insert the presentation content as context
                history_content.insert(0, types.Content(
                    role="user",
                    parts=[types.Part(text=f"Here is the content of the PowerPoint presentation:\n{full_ppt_text}")]
                ))

            except Exception as e:
                yield f"data: Error extracting .pptx: {str(e)}\n\n"
                return

        # Add the actual user prompt
        current_content = types.Content(
            role="user",
            parts=[types.Part(text=prompt)]
        )

        contents = history_content + [current_content]

        # Stream response from Gemini
        stream = genai_client.models.generate_content_stream(
            model="gemini-2.0-flash",
            contents=contents,
            config=types.GenerateContentConfig(
                system_instruction=system_prompt,
                max_output_tokens=max_tokens,
                temperature=temperature
            ),
        )

        for chunk in stream:
            yield f"{chunk.text}"

    except Exception as e:
        yield f"data: Error: {str(e)}\n\n"


async def ppt_openai_stream_response(
    prompt: str, 
    system_prompt: str, 
    max_tokens: int, 
    temperature: float,
    conversation_history: List[Dict[str, str]] = [],
    ppt_url: str = None
):
    try:
        input_data = []

        # Add conversation history
        for msg in conversation_history:
            input_data.append({
                "role": msg["role"],
                "content": [{"type": "input_text", "text": msg["content"]}]
            })

         # Extract PPTX content
        if ppt_url:
            try:
                response = httpx.get(ppt_url)
                response.raise_for_status()

                prs = Presentation(BytesIO(response.content))
                ppt_text = []

                for i, slide in enumerate(prs.slides):
                    slide_text = [f"--- Slide {i + 1} ---"]
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text.strip())
                    ppt_text.append("\n".join(slide_text))

                full_ppt_text = "\n\n".join(ppt_text)
                print(full_ppt_text)

                # Insert the presentation content as context
                ppt_input = {
                    "type": "input_text",
                    "text": f"Here is the content of the PowerPoint presentation:\n{full_ppt_text}"
                }

                user_content = [ppt_input]

            except Exception as e:
                yield f"data: Error extracting .pptx: {str(e)}\n\n"
                return

        else:
            user_content = []

        # Add user's actual prompt
        user_content.append({
            "type": "input_text",
            "text": prompt
        })

        # Append user role input to final input_data
        input_data.append({
            "role": "user",
            "content": user_content
        })

        # Stream response from OpenAI
        stream = openai_client.responses.create(
            model="gpt-4o",
            input=input_data,
            temperature=temperature,
            instructions=system_prompt,
            stream=True
        )

        for chunk in stream:
            if hasattr(chunk, "delta"):
                yield chunk.delta

    except Exception as e:
        yield f"data: Error: {str(e)}\n\n"